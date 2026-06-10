from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR_TYPE

# ─── Colors ──────────────────────────────────────────────────────────────────
def rgb(h): return RGBColor(int(h[:2],16), int(h[2:4],16), int(h[4:],16))

NAVY    = rgb('1E3A5F'); NAVYMID = rgb('2D5282'); NAVYLT  = rgb('3B6CB7')
TEAL    = rgb('0D9488'); TEALLT  = rgb('14B8A6'); TEALPAL = rgb('CCFBF1')
WHITE   = rgb('FFFFFF'); OFFWHT  = rgb('F8FAFC'); SLATE   = rgb('F1F5F9')
TEXTDK  = rgb('1E293B'); TEXTMD  = rgb('475569'); TEXTLT  = rgb('94A3B8')
AMBER   = rgb('F59E0B'); AMBERLT = rgb('FED7AA'); AMBERBG = rgb('FFF7ED'); AMBERDK = rgb('92400E')
BORDER  = rgb('CBD5E1'); BORDER2 = rgb('E2E8F0')
PURPLE  = rgb('7C3AED'); BROWN   = rgb('B45309'); BLUE    = rgb('0369A1')
# Terminal colors
T_BG    = rgb('0F172A')  # terminal background
T_GRAY  = rgb('94A3B8')  # normal output
T_WHITE = rgb('E2E8F0')  # important output
T_GREEN = rgb('4ADE80')  # >> prompt
T_YEL   = rgb('FCD34D')  # user input
T_CYAN  = rgb('38BDF8')  # menu items / headers
T_TEAL  = rgb('2DD4BF')  # separator lines

RECT=1; RRECT=5; OVAL=9

# ─── Helpers ─────────────────────────────────────────────────────────────────
def bg(sl, c):
    f = sl.background.fill; f.solid(); f.fore_color.rgb = c

def shp(sl, sid, x, y, w, h, fc, lc=None, lw=0.75, adj=None):
    s = sl.shapes.add_shape(sid, Inches(x), Inches(y), Inches(w), Inches(h))
    if fc: s.fill.solid(); s.fill.fore_color.rgb = fc
    else:  s.fill.background()
    if lc: s.line.color.rgb = lc; s.line.width = Pt(lw)
    else:  s.line.fill.background()
    if adj is not None and sid == RRECT:
        try: s.adjustments[0] = adj
        except: pass
    return s

def txt(sl, text, x, y, w, h, sz=12, bold=False, col=TEXTDK,
        al=PP_ALIGN.LEFT, ital=False, font='Calibri', va=MSO_ANCHOR.MIDDLE):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = va
    p = tf.paragraphs[0]; p.alignment = al
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = ital
    r.font.color.rgb = col; r.font.name = font
    return tb

def mtxt(sl, lines, x, y, w, h, sz=12, bold=False, col=TEXTDK,
         al=PP_ALIGN.LEFT, font='Calibri', ital=False):
    """lines: str | (str, bold) | (str, bold, color)"""
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(lines):
        if isinstance(line, str):   lb, lb2, lc = line, bold, col
        elif len(line) == 2:        lb, lb2, lc = line[0], line[1], col
        else:                       lb, lb2, lc = line
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al
        r = p.add_run(); r.text = lb
        r.font.size = Pt(sz); r.font.bold = lb2; r.font.italic = ital
        r.font.color.rgb = lc; r.font.name = font
    return tb

def ln(sl, x1, y1, x2, y2, col=TEXTLT, w=1):
    c = sl.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = col; c.line.width = Pt(w); return c

def slide_title(sl, text, bg_light=True):
    c = NAVY if bg_light else WHITE
    txt(sl, text, 0.5, 0.22, 10, 0.72, sz=32, bold=True, col=c,
        font='Cambria', va=MSO_ANCHOR.TOP)

# ─── Terminal helper ──────────────────────────────────────────────────────────
def terminal(sl, lines, x, y, w, h, sz=8.5):
    """
    lines: list of str | (str, color)
    Default color = T_GRAY
    """
    shp(sl, RRECT, x, y, w, h, T_BG, rgb('334155'), 0.8, adj=0.03)
    tb = sl.shapes.add_textbox(Inches(x+0.18), Inches(y+0.15),
                               Inches(w-0.36), Inches(h-0.3))
    tf = tb.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(lines):
        text, col = (item, T_GRAY) if isinstance(item, str) else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = text
        r.font.size = Pt(sz); r.font.color.rgb = col; r.font.name = 'Courier New'
    return tb

def callout_card(sl, x, y, w, h, header, body_lines, header_col, icon=''):
    shp(sl, RRECT, x, y, w, h, WHITE, BORDER2, 0.6, adj=0.07)
    txt(sl, icon + ' ' + header, x+0.15, y+0.1, w-0.3, 0.38,
        sz=12.5, bold=True, col=header_col)
    mtxt(sl, body_lines, x+0.15, y+0.52, w-0.3, h-0.62,
         sz=10.5, col=TEXTMD)

# ─────────────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BL = prs.slide_layouts[6]

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – 타이틀
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BL); bg(sl, NAVY)
shp(sl, OVAL, -1.5, -1.5, 6, 6, NAVYMID)
shp(sl, OVAL,  9.8,  4.0, 5, 5, NAVYMID)
txt(sl, '온라인 서점 시스템', 0.5, 2.05, 12.33, 1.4,
    sz=50, bold=True, col=WHITE, al=PP_ALIGN.CENTER, font='Cambria')
txt(sl, 'Python OOP 기반 콘솔 애플리케이션', 0.5, 3.65, 12.33, 0.7,
    sz=23, col=TEALLT, al=PP_ALIGN.CENTER)
tags = ['Python 3', 'OOP', '3계층 아키텍처', '딕셔너리 DB', '메뉴 패턴']
cw, cg = 2.0, 0.22
sx = (13.33 - (cw*len(tags) + cg*(len(tags)-1))) / 2
for i, t in enumerate(tags):
    cx = sx + i*(cw+cg)
    shp(sl, RRECT, cx, 4.75, cw, 0.52, TEAL, TEALLT, 0.8, adj=0.1)
    txt(sl, t, cx, 4.75, cw, 0.52, sz=13, col=WHITE, al=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – 프로젝트 개요
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BL); bg(sl, OFFWHT)
slide_title(sl, '프로젝트 개요')
stats = [
    ('5',    '주요 도메인',  '회원·도서·장바구니\n주문·배송'),
    ('3',    '계층 구조',    'Model · DAO\nService · App'),
    ('20+',  '유스케이스',   '회원·비회원\n관리자 기능 포함'),
    ('OOP',  'Python 기법',  '캡슐화·클래스변수\n__str__ 오버라이딩'),
]
cw, ch, cg = 2.85, 2.5, 0.33
sx = (13.33 - (cw*4 + cg*3)) / 2
for i, (n, label, sub) in enumerate(stats):
    cx = sx + i*(cw+cg)
    shp(sl, RRECT, cx, 1.2, cw, ch, WHITE, BORDER2, 0.7, adj=0.08)
    txt(sl, n, cx, 1.3, cw, 0.9, sz=44, bold=True, col=TEAL, al=PP_ALIGN.CENTER, font='Cambria')
    txt(sl, label, cx, 2.22, cw, 0.42, sz=15, bold=True, col=TEXTDK, al=PP_ALIGN.CENTER)
    txt(sl, sub,   cx, 2.65, cw, 0.9,  sz=11, col=TEXTMD, al=PP_ALIGN.CENTER)
shp(sl, RRECT, 0.4, 4.0, 12.53, 3.1, TEALPAL, TEALLT, 1, adj=0.05)
txt(sl, '▍ 프로젝트 소개', 0.7, 4.1, 5, 0.45, sz=15, bold=True, col=NAVY)
mtxt(sl, [
    '●  회원(비회원 / 일반회원 / 관리자)이 이용하는 Python 콘솔 기반 온라인 서점',
    '●  도서 조회 → 장바구니 → 주문 → 배송까지 완전한 구매 흐름 구현',
    '●  파일·DB 없이 딕셔너리 기반 인메모리 저장소 사용 (3계층 아키텍처)',
    '●  Python OOP 핵심 기법: 캡슐화(name mangling), 클래스 변수, __str__ 오버라이딩',
], 0.7, 4.62, 12.0, 2.3, sz=14, col=TEXTDK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – 시스템 아키텍처
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BL); bg(sl, OFFWHT)
slide_title(sl, '시스템 아키텍처 (3계층 구조)')
layers = [
    (NAVY,        '1. Presentation Layer', 'BookShopApp',
     'select_menu() · run_xxx_menu() · menu_xxx() 메뉴 패턴 적용'),
    (TEAL,        '2. Service Layer', 'MemberService · BookService · CartService · OrderService · DeliveryService',
     '비즈니스 로직 처리, DAO 조합, 데이터 검증'),
    (NAVYLT,      '3. DAO Layer', 'MemberDAO · BookDAO · CartDAO · OrderDAO · DeliveryDAO',
     '딕셔너리 기반 CRUD, auto-increment __next_no 관리'),
    (rgb('334155'), '4. Model Layer', 'Member · Book · Cart · Order · Delivery',
     '도메인 데이터 보관, name mangling 캡슐화, __str__ 오버라이딩'),
]
lx, lw, lh, lg = 1.5, 10.3, 1.2, 0.22
for i, (color, title, names, desc) in enumerate(layers):
    ly = 1.15 + i*(lh+lg)
    shp(sl, RRECT, lx, ly, lw, lh, color, None, adj=0.06)
    txt(sl, title, lx+0.2, ly+0.05, 3.2, 0.45, sz=14, bold=True, col=WHITE)
    txt(sl, names, lx+3.5, ly+0.04, 4.5, 0.45, sz=11.5, bold=True, col=rgb('A5F3FC'))
    txt(sl, desc,  lx+0.2, ly+0.55, lw-0.4, 0.55, sz=10.5, col=rgb('CBD5E1'))
    if i < 3:
        ay = ly + lh
        ln(sl, 6.65, ay, 6.65, ay+lg, col=TEALLT, w=1.5)
txt(sl, '각 계층은 단방향 의존성 — 상위 계층이 하위 계층을 호출합니다',
    1.5, 6.65, 10.3, 0.55, sz=11, col=TEXTLT, al=PP_ALIGN.CENTER, ital=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – 유스케이스 다이어그램
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BL); bg(sl, WHITE)
slide_title(sl, '유스케이스 다이어그램', bg_light=True)
shp(sl, RRECT, 1.9, 1.05, 11.0, 6.15, SLATE, BORDER, 1.2, adj=0.04)
txt(sl, '온라인 서점 시스템', 5.5, 1.08, 4, 0.35, sz=10.5, col=TEXTLT,
    al=PP_ALIGN.CENTER, ital=True)
group_bands = [
    (1.95, 1.45, 10.9, 1.55, rgb('EFF6FF'), rgb('BFDBFE')),
    (1.95, 3.05, 10.9, 2.3,  rgb('F0FDF4'), rgb('86EFAC')),
    (1.95, 5.4,  10.9, 1.7,  rgb('FDF4FF'), rgb('D8B4FE')),
]
for (bx, by, bw, bh, fc, lc) in group_bands:
    shp(sl, RRECT, bx, by, bw, bh, fc, lc, 0.7, adj=0.04)
actors_data = [
    ('비회원', 1.58, rgb('2563EB'), WHITE,
     ['도서 목록 보기', '회원가입', '로그인'], [1.68]),
    ('회원',   3.55, TEAL, WHITE,
     ['장바구니 관리', '주문하기', '주문/배송 조회', '내 정보 관리', '로그아웃'], [3.18, 3.75]),
    ('관리자', 5.78, PURPLE, WHITE,
     ['도서 관리', '회원 관리', '주문 관리', '배송 관리'], [5.52]),
]
for (aname, ay, acol, atcol, ucs, uc_rows) in actors_data:
    shp(sl, RRECT, 0.1, ay, 1.65, 0.48, acol, None, adj=0.1)
    txt(sl, aname, 0.1, ay, 1.65, 0.48, sz=14, bold=True, col=atcol, al=PP_ALIGN.CENTER)
    uc_idx, uc_w = 0, 3.0
    for row_y in uc_rows:
        row_ucs = ucs[uc_idx:uc_idx+3]; uc_idx += 3
        n = len(row_ucs)
        sx2 = 2.05 + (10.9 - (n*uc_w + (n-1)*0.22)) / 2
        for j, uc in enumerate(row_ucs):
            ux = sx2 + j*(uc_w+0.22)
            shp(sl, OVAL, ux, row_y, uc_w, 0.45, WHITE, acol, 1.2)
            txt(sl, uc, ux, row_y, uc_w, 0.45, sz=11, col=TEXTDK, al=PP_ALIGN.CENTER)
    ln(sl, 1.75, ay+0.24, 2.05, uc_rows[0]+0.22, col=acol, w=0.8)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – 클래스 다이어그램
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BL); bg(sl, OFFWHT)
slide_title(sl, '클래스 다이어그램 (모델 계층)')
classes = [
    (TEAL,   'Member',
     ['member_no: int', 'id: str', 'password: str', 'name: str', 'phone: str', 'address: str'],
     ['get/set_*()', 'update_info(name,phone,addr)', 'check_password(pw)']),
    (NAVY,   'Book',
     ['book_no: int', 'title: str', 'author: str', 'publisher: str', 'price: int', 'stock: int'],
     ['get/set_*()', 'update_info()', 'decrease_stock(qty)', 'increase_stock(qty)']),
    (PURPLE, 'Cart',
     ['member_no: int', 'items: dict', '  Key: book_no', '  Value: qty'],
     ['add_book(book_no,qty)', 'change_quantity(no,qty)', 'remove_book(book_no)', 'clear_cart()']),
    (BROWN,  'Order',
     ['order_no: int', 'member_no: int', 'order_items: dict', 'total_price: int', 'order_status: str'],
     ['get_*()', 'change_status(status)', 'STATUS 상수 5종']),
    (BLUE,   'Delivery',
     ['order_no: int', 'member_no: int', 'address: str', 'status: str'],
     ['get_*()', '※ View-Only 모델', '(OrderDAO 기반 생성)']),
]
n = len(classes)
cw = (13.33 - 0.3 - (n-1)*0.18) / n
cx0 = 0.15
for i, (color, name, attrs, methods) in enumerate(classes):
    cx = cx0 + i*(cw+0.18)
    shp(sl, RRECT, cx, 1.05, cw, 0.5, color, None, adj=0.06)
    txt(sl, name, cx, 1.05, cw, 0.5, sz=14, bold=True, col=WHITE, al=PP_ALIGN.CENTER)
    attr_h = 0.3 * len(attrs) + 0.15
    shp(sl, RECT, cx, 1.55, cw, attr_h, WHITE, BORDER, 0.5)
    mtxt(sl, attrs, cx+0.08, 1.6, cw-0.16, attr_h-0.1, sz=9, col=TEXTDK)
    meth_y = 1.55 + attr_h
    meth_h = 0.3 * len(methods) + 0.15
    shp(sl, RECT, cx, meth_y, cw, meth_h, SLATE, BORDER, 0.5)
    mtxt(sl, methods, cx+0.08, meth_y+0.05, cw-0.16, meth_h-0.1, sz=9, col=TEXTMD, ital=True)
txt(sl, '각 모델마다 DAO(저장소)와 Service(비즈니스 로직) 클래스가 존재합니다 — BookDAO / BookService, OrderDAO / OrderService …',
    0.3, 7.1, 12.7, 0.35, sz=10, col=TEXTLT, al=PP_ALIGN.CENTER, ital=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – 주문 프로세스
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BL); bg(sl, WHITE)
slide_title(sl, '주요 기능 흐름 – 주문 프로세스')
steps = [
    ('1', '로그인',     '아이디·비밀번호 확인\n장바구니 자동 생성'),
    ('2', '도서 선택',  '도서 목록 조회\n장바구니 담기\n(재고 실시간 확인)'),
    ('3', '주문 확인',  '전화번호·주소 없으면\n주문 전 입력 요청'),
    ('4', '주문 완료',  '재고 즉시 차감\n장바구니 자동 초기화'),
    ('5', '배송 관리',  '관리자가 상태 변경\n주문접수 → 배송중\n→ 배송완료'),
]
n_s = len(steps)
cw_s, ch_s = 2.15, 3.1
cg_s = (13.33 - 0.4 - n_s*cw_s) / (n_s-1)
sx3 = 0.2
for i, (num, title, desc) in enumerate(steps):
    cx = sx3 + i*(cw_s+cg_s)
    shp(sl, RRECT, cx, 1.1, cw_s, ch_s, WHITE, BORDER2, 0.7, adj=0.08)
    shp(sl, OVAL, cx+(cw_s-0.6)/2, 1.2, 0.6, 0.6, TEAL)
    txt(sl, num, cx+(cw_s-0.6)/2, 1.2, 0.6, 0.6, sz=16, bold=True, col=WHITE, al=PP_ALIGN.CENTER)
    txt(sl, title, cx, 1.92, cw_s, 0.48, sz=14, bold=True, col=NAVY, al=PP_ALIGN.CENTER)
    txt(sl, desc,  cx+0.1, 2.45, cw_s-0.2, 1.6, sz=11, col=TEXTMD, al=PP_ALIGN.CENTER)
    if i < n_s-1:
        shp(sl, RECT, cx+cw_s, 2.5, cg_s, 0.06, TEALLT)
shp(sl, RRECT, 0.4, 4.55, 12.53, 1.65, AMBERBG, AMBERLT, 1, adj=0.05)
txt(sl, '⚠  예외 처리', 0.7, 4.65, 3, 0.42, sz=14, bold=True, col=BROWN)
mtxt(sl, [
    '●  장바구니 비어있으면 → 도서를 직접 선택해 바로 주문 가능',
    '●  주문 취소 상태 → 관리자가 주문 거부 불가',
    '●  취소·거부 상태 → 배송 상태 수정 불가  |  장바구니 초과 담기 → 재고 기준으로 차단',
], 0.7, 5.12, 12.0, 1.0, sz=13, col=AMBERDK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – 동작 화면 ① 시작 & 로그인
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BL); bg(sl, rgb('0F172A'))
txt(sl, '동작 화면 ①  —  시작 화면 & 로그인', 0.4, 0.2, 12.5, 0.65,
    sz=26, bold=True, col=WHITE, font='Cambria', va=MSO_ANCHOR.TOP)

terminal(sl, [
    ('======== Online Book Store ==========',          T_CYAN),
    ('>>>>도서 목록<<<<',                               T_CYAN),
    ('=' * 62,                                          T_TEAL),
    (' 번호 |  제목                          |  저자         |   가격  | 재고', T_WHITE),
    ('=' * 62,                                          T_TEAL),
    ('   1  |  벨 푸페의 슈퍼달링 약혼       |  아사기리 아사키 | 10000원 | 15권', T_GRAY),
    ('   2  |  패배 히로인이 너무 많아!      |  아마모리 타키비 |  9000원 | 12권', T_GRAY),
    ('   3  |  Re: 제로부터 시작하는 이세계  |  나가츠키 탓페이 |  7200원 | 10권', T_GRAY),
    ('   4  |  마녀의 여행                   |  시라이시 조우기 |  9500원 |  8권', T_GRAY),
    ('   5  |  전생했더니 슬라임이었던 건    |  후세           | 10000원 |  5권', T_GRAY),
    ('   6  |  사일런트 위치                 |  이소라 마츠리  |  9000원 |  7권', T_GRAY),
    ('=' * 62,                                          T_TEAL),
    ('',                                                T_GRAY),
    ('=' * 62,                                          T_TEAL),
    ('1. 로그인   2. 회원가입   0. 종료',               T_CYAN),
    ('=' * 62,                                          T_TEAL),
    ('>> 메뉴 : 1',                                     T_GREEN),
    ('',                                                T_GRAY),
    ('>>>>>>>> 로그인 <<<<<<<<<',                       T_CYAN),
    ('아이디 : user1',                                  T_YEL),
    ('비밀번호 : ****',                                  T_YEL),
    ('이유찬님, 환영합니다.',                            T_WHITE),
], x=0.25, y=1.0, w=7.9, h=6.3, sz=8.5)

# 오른쪽 설명 카드
callout_card(sl, 8.35, 1.05, 4.75, 1.8, '실행 즉시 도서 목록 표시',
    ['로그인 없이도 도서 목록을\n먼저 보여줌 (비회원 접근)'],
    TEAL, '📋')
callout_card(sl, 8.35, 2.95, 4.75, 1.9, '한글 너비 정렬',
    ['str_width() 로 한글(2칸)·영문(1칸)\n너비를 계산해 테이블 정렬\nljust_k() · rjust_k() 사용'],
    NAVYLT, '📐')
callout_card(sl, 8.35, 4.95, 4.75, 1.8, '로그인 성공',
    ['회원 조회 후 환영 메시지 출력\n자동으로 장바구니(Cart) 객체 생성\n관리자면 관리자 메뉴로 분기'],
    TEAL, '🔑')

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – 동작 화면 ② 장바구니 & 주문
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BL); bg(sl, rgb('0F172A'))
txt(sl, '동작 화면 ②  —  장바구니 & 주문', 0.4, 0.2, 12.5, 0.65,
    sz=26, bold=True, col=WHITE, font='Cambria', va=MSO_ANCHOR.TOP)

terminal(sl, [
    ('1. 장바구니 관리  2. 주문하기  3. 주문/배송 조회  4. 내 정보 관리  0. 로그아웃', T_CYAN),
    ('>> 메뉴 : 1',                         T_GREEN),
    ('>>>>장바구니 관리<<<<',                T_CYAN),
    ('1. 장바구니 담기  2. 장바구니 보기  3. 장바구니 삭제  4. 주문하기  0. 뒤로가기', T_CYAN),
    ('>> 메뉴 : 1',                         T_GREEN),
    ('>>>>장바구니 담기<<<<',                T_CYAN),
    ('>> 도서 번호 : 1',                    T_YEL),
    ('>> 수량 : 2',                         T_YEL),
    ('장바구니에 담았습니다',               T_WHITE),
    ('',                                    T_GRAY),
    ('>> 메뉴 : 2',                         T_GREEN),
    ('>>>>장바구니 보기<<<<',                T_CYAN),
    ('[1] 벨 푸페의 슈퍼달링 약혼 x2 = 20000원', T_WHITE),
    ('합계 : 20000원',                      T_WHITE),
    ('',                                    T_GRAY),
    ('>> 메뉴 : 4',                         T_GREEN),
    ('>>>>주문하기<<<<',                     T_CYAN),
    ('벨 푸페의 슈퍼달링 약혼 x2',          T_WHITE),
    ('총 결제 금액 : 20000원',              T_WHITE),
    ('주문하시겠습니까? (y/n) : y',         T_YEL),
    ('주문 완료',                           rgb('4ADE80')),
    ('',                                    T_GRAY),
    ('>>> 재고 자동 차감: 15권 → 13권',     rgb('FCD34D')),
    ('>>> 장바구니 자동 초기화 완료',        rgb('FCD34D')),
], x=0.25, y=1.0, w=7.9, h=6.3, sz=8.5)

callout_card(sl, 8.35, 1.05, 4.75, 1.75, '장바구니 메뉴 구조',
    ['담기·보기·삭제·주문하기가\n하나의 서브메뉴로 그룹화\n장바구니 안에서도 주문 가능'],
    TEAL, '🛒')
callout_card(sl, 8.35, 2.9, 4.75, 1.75, '재고 초과 방지',
    ['장바구니에 담긴 수량 + 추가 수량이\n재고를 초과하면 즉시 차단\nalready_in_cart 합산 체크'],
    NAVYLT, '🔒')
callout_card(sl, 8.35, 4.75, 4.75, 2.0, '주문 완료 처리',
    ['y 확인 후 3가지 동시 처리:\n① 재고 즉시 차감 (decrease_stock)\n② 주문 객체 생성 (create_order)\n③ 장바구니 초기화 (clear_cart)'],
    BROWN, '📦')

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – 동작 화면 ③ 관리자
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BL); bg(sl, rgb('0F172A'))
txt(sl, '동작 화면 ③  —  관리자 화면', 0.4, 0.2, 12.5, 0.65,
    sz=26, bold=True, col=WHITE, font='Cambria', va=MSO_ANCHOR.TOP)

terminal(sl, [
    ('>>>>>>>> 로그인 <<<<<<<<<',            T_CYAN),
    ('아이디 : admin',                       T_YEL),
    ('비밀번호 : *****',                      T_YEL),
    ('관리자님, 환영합니다.',                 T_WHITE),
    ('',                                     T_GRAY),
    ('1. 도서 관리 메뉴  2. 회원 관리 메뉴  3. 주문 관리 메뉴  4. 배송 관리 메뉴  0. 로그아웃', T_CYAN),
    ('>> 메뉴 : 1',                          T_GREEN),
    ('>>>>도서 관리 메뉴<<<<',               T_CYAN),
    ('1. 도서 목록  2. 도서 등록  3. 도서 수정  4. 도서 삭제  0. 뒤로가기', T_CYAN),
    ('>> 메뉴 : 3',                          T_GREEN),
    ('>>>>도서 수정<<<<',                     T_CYAN),
    ('>> 수정할 도서 번호 : 1',              T_YEL),
    ('1. 제목 변경  2. 저자 변경  3. 출판사 변경  4. 가격 변경  5. 재고 변경  0. 취소', T_CYAN),
    ('>> 메뉴 : 5',                          T_GREEN),
    ('>>>>재고 변경<<<<',                     T_CYAN),
    ('>> 새 재고 [15] : 20',                 T_YEL),
    ('재고 변경 성공',                        rgb('4ADE80')),
    ('',                                     T_GRAY),
    ('[ 배송 관리 예시 ]',                   T_TEAL),
    ('주문번호: 1  회원: 이유찬  주소: 지옥  상태: 주문접수', T_GRAY),
    ('>> 수정할 주문 번호 : 1',              T_YEL),
    ('1. 배송중   2. 배송완료',              T_CYAN),
    ('>> 상태 선택 : 1',                     T_YEL),
    ('수정 성공  →  상태: 배송중',           rgb('4ADE80')),
], x=0.25, y=1.0, w=7.9, h=6.3, sz=8.5)

callout_card(sl, 8.35, 1.05, 4.75, 1.75, '관리자 전용 메뉴',
    ['로그인 시 admin ID 감지 →\n자동으로 관리자 메뉴로 분기\n일반 유저 메뉴와 완전 분리'],
    PURPLE, '🔐')
callout_card(sl, 8.35, 2.9, 4.75, 1.8, '도서 수정 메뉴 방식',
    ['도서 선택 후 항목(제목/저자/\n출판사/가격/재고)을 메뉴로 선택\n선택한 항목만 수정 후 자동 복귀'],
    NAVY, '✏️')
callout_card(sl, 8.35, 4.8, 4.75, 1.95, '배송 상태 제어',
    ['취소·거부 주문은 배송 상태\n수정 불가 (자동 차단)\n주문접수 → 배송중 → 배송완료\n단방향 흐름 보장'],
    TEAL, '🚚')

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – Python 기술 포인트
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BL); bg(sl, OFFWHT)
slide_title(sl, 'Python 기술 포인트')
points = [
    ('🔒', '캡슐화 (Name Mangling)',
     '변수명에 __ 접두사로 외부 접근 차단\n→ get/set 메서드로만 데이터 접근 가능'),
    ('📋', '클래스 변수 메뉴 패턴',
     '메뉴 목록을 클래스 변수로 선언\n→ select_menu() + run_xxx_menu() 구조'),
    ('📖', '__str__ 오버라이딩',
     'print(객체) 호출 시\n→ 가독성 있는 포맷으로 자동 출력'),
    ('🗂️', '딕셔너리 기반 저장소',
     '파일·DB 없이 메모리 딕셔너리 CRUD\n→ Key: auto-increment 번호'),
    ("↩️", "'backhome' 패턴",
     "return 'backhome' 으로\n→ 다중 중첩 메뉴 즉시 탈출"),
    ('📐', '한글 출력 정렬',
     'str_width() · ljust_k() 커스텀 함수\n→ 한글 2칸 너비 고려한 테이블 정렬'),
]
cols, rows = 3, 2
cw2, ch2, gx, gy = 3.8, 2.2, 0.35, 0.3
sx4 = (13.33 - (cols*cw2 + (cols-1)*gx)) / 2
for i, (icon, title, desc) in enumerate(points):
    col2 = i % cols; row2 = i // cols
    cx = sx4 + col2*(cw2+gx); cy = 1.1 + row2*(ch2+gy)
    shp(sl, RRECT, cx, cy, cw2, ch2, WHITE, BORDER2, 0.6, adj=0.07)
    txt(sl, icon + '  ' + title, cx+0.18, cy+0.15, cw2-0.36, 0.5, sz=13, bold=True, col=NAVY)
    txt(sl, desc, cx+0.18, cy+0.7, cw2-0.36, 1.35, sz=11.5, col=TEXTMD)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – 마무리
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BL); bg(sl, NAVY)
shp(sl, OVAL, -1.5, -1.5, 6, 6, NAVYMID)
shp(sl, OVAL, 10.0,  4.5, 5, 5, NAVYMID)
txt(sl, '감사합니다', 0.5, 2.0, 12.33, 1.4,
    sz=54, bold=True, col=WHITE, al=PP_ALIGN.CENTER, font='Cambria')
txt(sl, 'Python OOP 기반 온라인 서점 시스템', 0.5, 3.6, 12.33, 0.7,
    sz=22, col=TEALLT, al=PP_ALIGN.CENTER)
summary = ['5개 도메인', '3계층 아키텍처', '20+ 유스케이스', 'Python OOP 집약']
sw, sg = 2.3, 0.25
ssx = (13.33 - (sw*len(summary) + sg*(len(summary)-1))) / 2
for i, t in enumerate(summary):
    scx = ssx + i*(sw+sg)
    shp(sl, RRECT, scx, 4.6, sw, 0.55, TEAL, TEALLT, 0.8, adj=0.1)
    txt(sl, t, scx, 4.6, sw, 0.55, sz=14, col=WHITE, al=PP_ALIGN.CENTER)

# ─── 저장 ────────────────────────────────────────────────────────────────────
out = r'C:\Lecture\python26\과제\OnlineBookStore\OnlineBookStore_PPT.pptx'
prs.save(out)
print('PPT saved:', out)
